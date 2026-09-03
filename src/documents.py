"""
documents_tab.py
================
Streamlit "Documents" tab for the evaluation pipeline.
Handles raw document ingestion (file upload, JSON, CSV, SQL)
and preprocessed chunk ingestion (JSON, CSV, SQL).

Depends on:
    pip install streamlit pandas python-docx pypdf striprtf
                ebooklib beautifulsoup4 pyodbc sqlalchemy
                markdown openpyxl python-pptx
"""

import io
import json
import textwrap
from pathlib import Path
import hashlib
from nlp_analytics_tab import render_nlp_analytics_tab

import pandas as pd
import streamlit as st
import puremagic
import itertools
from chonkie import (
    CodeChunker,
    FastChunker,
    LateChunker,
    NeuralChunker,
    RecursiveChunker,
    SemanticChunker,
    SentenceChunker,
    TableChunker,
    TokenChunker
)
from collections import Counter
import asyncio
import queue
from background import (
    submit,
    cancel,
    get_state,
    set_state,
    delete_value,
    close_dialog,
    submit_coroutines,
    fragment
)
from templates import metric_card
import os
import copy
import re
import time

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

SUPPORTED_DOC_TYPES = [
    # plain text
    "txt", "md", "rst", "log",
    # rich text
    "rtf", "odt",
    # Microsoft Office
    "docx", "doc", "pptx", "ppt", "xlsx", "xls",
    # PDF
    "pdf",
    # web / markup
    "html", "htm", "xml",
    # e-book
    "epub",
    # structured / code
    "json", "yaml", "yml", "toml", "csv",
]

CONTENT_KEY   = "content"   # expected key in JSON / column in CSV

# ---------------------------------------------------------------------------
# Parsers — each returns a list of {"content": str, **metadata} dicts
# ---------------------------------------------------------------------------

def parser(parse_func):
    def parser_wrapper(input):
        try:
            data = parse_func(input)
            return {'content': [data] if isinstance(data, str) else data}
        except BaseException as _:
            return None
    return parser_wrapper

@parser
def _parse_txt_md(file_bytes: bytes):
    return file_bytes.decode("utf-8", errors="replace")

@parser
def _parse_pdf(file_bytes: bytes):
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(file_bytes))
    text = "\n".join(p.extract_text() or "" for p in reader.pages)
    return text

@parser
def _parse_docx(file_bytes: bytes):
    import docx
    doc  = docx.Document(io.BytesIO(file_bytes))
    text = "\n".join(p.text for p in doc.paragraphs)
    return text

@parser
def _parse_rtf(file_bytes: bytes):
    from striprtf.striprtf import rtf_to_text
    text = rtf_to_text(file_bytes.decode("utf-8", errors="replace"))
    return text

@parser
def _parse_epub(file_bytes: bytes):
    import ebooklib
    from ebooklib import epub
    from bs4 import BeautifulSoup
    book  = epub.read_epub(io.BytesIO(file_bytes))
    parts = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_content(), "html.parser")
        parts.append(soup.get_text())
    return "\n".join(parts)

@parser
def _parse_html(file_bytes: bytes):
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(file_bytes, "html.parser")
    return soup.get_text()

@parser
def _parse_pptx(file_bytes: bytes, file_name: str):
    from pptx import Presentation
    prs  = Presentation(io.BytesIO(file_bytes))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

@parser
def _parse_xlsx(file_bytes: bytes, file_name: str):
    df   = pd.read_excel(io.BytesIO(file_bytes))
    text = df.to_string(index=False)
    return text

PARSER_MAP = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": _parse_docx,
    "application/msword": _parse_docx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": _parse_pptx,
    "application/vnd.ms-powerpoint": _parse_pptx,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": _parse_xlsx,
    "application/vnd.ms-excel": _parse_xlsx,
    "application/pdf": _parse_pdf,
    "application/epub+zip": _parse_epub,
    "application/rtf": _parse_rtf,
    "text/html": _parse_html,
    "text/xml": _parse_txt_md,
    "text/plain": _parse_txt_md
}

def get_mime(file_bytes: bytes) -> str:
    try:
        matches = puremagic.magic_string(file_bytes)
        return matches[0].mime_type if matches else "application/octet-stream"
    except BaseException as _:
        return "application/octet-stream"

def parse_uploaded_file(uploaded_file) -> list[dict]:
    """Dispatch to the right parser based on true file type."""
    file_bytes = uploaded_file.read()
    file_type = get_mime(file_bytes)
    parser = PARSER_MAP.get(file_type, _parse_txt_md)
    out = parser(file_bytes)
    return out

@parser
def _validate_content_key(records: list[dict]) -> list[dict]:
    """Filter out records missing the 'content' key and warn the user."""
    valid = [r[CONTENT_KEY] for r in records if r.get(CONTENT_KEY)]
    if not valid:
        raise KeyError
    return valid

async def _records_from_json(file, queue, mode) -> list[dict]:
    raw = await asyncio.to_thread(json.load, file)
    if not isinstance(raw, list):
        queue.put((file.name, mode, None))
        return
    queue.put((file.name, mode, _validate_content_key(raw)))

async def _records_from_csv(file, queue, mode) -> list[dict]:
    df = await asyncio.to_thread(pd.read_csv, file)
    if CONTENT_KEY not in df.columns:
        queue.put((file.name, mode, None))
        return
    records = _validate_content_key(df.to_dict(orient="records"))
    queue.put((file.name, mode, records))

def _build_sql_engine(db_type, server, database, username, password, use_windows_auth):
    try:
        from sqlalchemy import create_engine
        if db_type == "SQL Server":
            if use_windows_auth:
                conn_str = (
                    f"mssql+pyodbc://{server}/{database}"
                    "?driver=ODBC+Driver+17+for+SQL+Server&trusted_connection=yes"
                )
            else:
                conn_str = (
                    f"mssql+pyodbc://{username}:{password}@{server}/{database}"
                    "?driver=ODBC+Driver+17+for+SQL+Server"
                )
        elif db_type == "PostgreSQL":
            conn_str = (
                f"postgresql+psycopg2://{username}:{password}@{server}/{database}"
            )
        return create_engine(conn_str)
    except BaseException as _:
        import traceback
        print("--- DATABASE ERROR DIRECTION ---")
        traceback.print_exc()
        return None

async def read_sql(query, engine, queue, filename, mode):
    try:
        if not engine:
            print('no engine')
            raise ValueError
        
        df = await asyncio.to_thread(pd.read_sql, query, engine)

        if CONTENT_KEY not in df.columns:
            raise KeyError
        
        records = _validate_content_key(df.to_dict(orient="records"))
        queue.put((filename, mode, records))
    
    except BaseException as _:
        import traceback
        print("--- DATABASE ERROR DIRECTION ---")
        traceback.print_exc()
        queue.put((
            filename,
            mode,
            None
        ))

def _sql_form(key_prefix: str, queue: queue.Queue(), mode):
    """Renders the SQL connection form. Returns (engine, query) or (None, None)."""
    db_type = st.selectbox("Database Type", ["PostgreSQL", "SQL Server"])
    col1, col2 = st.columns(2)
    server = col1.text_input("Server / Host", key=f"{key_prefix}_server")
    database = col2.text_input("Database", key=f"{key_prefix}_db")
    if db_type == "SQL Server":
        use_win = st.checkbox("Use Windows Authentication", key=f"{key_prefix}_winauth")
    else:
        use_win = False

    username = password = ""
    if not use_win:
        col3, col4 = st.columns(2)
        username = col3.text_input("Username", key=f"{key_prefix}_user")
        password = col4.text_input("Password", type="password", key=f"{key_prefix}_pass")
    
    col5, col6, col7 = st.columns(3)
    schema = col5.text_input("Schema (Optional)", key=f"{key_prefix}_schema")
    table = col6.text_input("Table Name", key=f"{key_prefix}_table")
    column = col7.text_input(f"Content Column", key=f"{key_prefix}_column")

    full_table_name = f"{schema + '.' if schema else ''}{table}"

    query = f"SELECT {column} FROM {full_table_name}"

    if st.button("Test & Load", key=f"{key_prefix}_connect"):
        engine = _build_sql_engine(db_type, server, database, username, password, use_win)
        filename = f'{server}.{database}.{full_table_name}.{column}'
        get_state('coro_queue').put(('uploading', read_sql(query, engine, queue, filename, mode)))
        set_state('active_dialog', value=None)
        st.rerun()

UPLOADED_CHUNK_SOURCES = ['application/x-sql-dataset', 'text/plain']

# ---------------------------------------------------------------------------
# Chunk inspector
# ---------------------------------------------------------------------------

import streamlit as st
import pandas as pd
import textwrap

# Ensure this matches your actual key name (e.g., "chunks" or "text")
CONTENT_KEY = "chunks" 

def _chunk_inspector(records: dict): # Changed type hint to dict based on your code
    if not records:
        st.info("No records to inspect.")
        return
    
    rows = len(get_state('chunks'))
    st.subheader("Chunk Inspector")
    idx = st.number_input("Expand chunk #", min_value=0,
                          max_value=rows - 1, step=1)
    with st.container(height=400, border=False):
        with st.expander(f"Chunk {idx} — full text", expanded=True):
            text = _load_txt(get_state('project_dir') / 'chunks' / f'{list(get_state('chunks').keys())[idx]}.txt')
            st.text(text)

# ---------------------------------------------------------------------------
# Chunking configuration + execution
# ---------------------------------------------------------------------------

@st.dialog("Chunking Settings")
def chunking_dialog(docs = list[dict]):
    chunkers = {
        "Code Chunker": CodeChunker,
        "Fast Chunker": FastChunker,
        "Late Chunker": LateChunker,
        "Neural Chunker": NeuralChunker,
        "Recursive Chunker": RecursiveChunker,
        "Semantic Chunker": SemanticChunker,
        "Sentence Chunker": SentenceChunker,
        "Table Chunker": TableChunker,
        "Token Chunker": TokenChunker
    }
    strategy = st.selectbox("Strategy", list(chunkers.keys()))
    chunking_method = chunkers[strategy]

    args = dict()

    if strategy == "Code Chunker":
        args['chunk_size'] = st.number_input("Max. Tokens per Chunk", value=2048, step=1, min_value=1)
    
    elif strategy == "Fast Chunker":
        args['chunk_size'] = st.number_input("Target Chunk Size (Bytes)", value=8192, step=1, min_value=64)
    
    elif strategy == "Late Chunker":
        args['chunk_size'] = st.number_input("Max. Tokens per Chunk", value=2048, step=1, min_value=1)
        args['min_characters_per_chunk'] = st.number_input("Min. Characters per Chunk", value=128, step=1, min_value=1)
    
    elif strategy == "Neural Chunker":
        args['min_characters_per_chunk'] = st.number_input("Min. Characters per Chunk", value=128, step=1, min_value=1)
    
    elif strategy == "Recursive Chunker":
        args['chunk_size'] = st.number_input("Max. Tokens per Chunk", value=2048, step=1, min_value=1)
        args['min_characters_per_chunk'] = st.number_input("Min. Characters per Chunk", value=128, step=1, min_value=1)
    
    elif strategy == "Semantic Chunker":
        args['threshold'] = st.number_input("Threshold", value=0.8, step=0.01, min_value=0.0)
        args['chunk_size'] = st.number_input("Max. Tokens per Chunk", value=2048, step=1, min_value=1)
        args['similarity_window'] = st.number_input("Similarity Window (Sentences)", value=3, step=1, min_value=1)
        args['min_sentences_per_chunk'] = st.number_input("Min. Sentences per Chunk", value=1, step=1, min_value=1)
        args['skip_window'] = st.number_input("Skip Window (Sentences)", value=0, step=1, min_value=0)
    
    elif strategy == "Sentence Chunker":
        args['chunk_size'] = st.number_input("Max. Tokens per Chunk", value=2048, step=1, min_value=1)
        args['chunk_overlap'] = st.number_input("Overlapping Tokens per Chunk", value=0, step=1, min_value=0)
        args['min_sentences_per_chunk'] = st.number_input("Min. Sentences per Chunk", value=1, step=1, min_value=1)
    
    elif strategy == "Table Chunker":
        args['chunk_size'] = st.number_input("Max. Rows per Chunk", value=3, step=1, min_value=1)
    
    elif strategy == "Token Chunker":
        args['chunk_size'] = st.number_input("Max. Tokens per Chunk", value=2048, step=1, min_value=1)
        args['chunk_overlap'] = st.number_input("% Overlap between Chunks", value=0.0, step=0.01, min_value=0.0)

    chunks_queue = get_state('chunk_queue')

    if st.button("Apply & Chunk"):
        get_state('coro_queue').put(('chunking', _chunk_documents(chunking_method, args, docs, chunks_queue)))
        set_state('active_dialog', value=None)
        st.rerun()

async def _chunk_documents(chunking_method, args, docs, queue) -> list[dict]:
    
    chunker = chunking_method(**args)
    chunker_config = {"class": type(chunker).__name__, **chunker.__dict__}
    chunker_string = json.dumps(chunker_config, sort_keys=True, default=str)

    chunker_hash = hash(chunker_string)

    for key in docs:

        if chunker_hash == docs[key].get('chunking_strategy'):
            continue

        else:
            try:
                doc_path = get_state('project_dir') / 'docs' / f'{key}.txt'
                doc = await asyncio.to_thread(_load_txt, doc_path)

                chunks = await asyncio.to_thread(chunker.chunk, doc)
                chunk_texts = {chunk.text for chunk in chunks}

                queue.put((key, chunker_hash, chunk_texts))
            except BaseException as _:
                queue.put((key, None, None))

def _load_txt(path):
    with open(path, 'r', encoding='utf-8') as f:
        txt = f.read()
    return txt

# ---------------------------------------------------------------------------
# Main tab renderer
# ---------------------------------------------------------------------------

def hash(text) -> str:
    return hashlib.sha256(text.encode()).hexdigest()

async def load_json(path: Path, queue: queue.Queue(), *attr):
    if path.exists():
        queue.put((*attr, await asyncio.to_thread(_load_json, path)))
    else:
        queue.put((*attr, dict()))

def _load_json(path):
    with open(path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    unserialize(data)
    return data

def unserialize(data: dict):
    for key in data:
        if key in {'docs', 'chunks'}:
            data[key] = set(data[key])
        elif isinstance(data[key], dict):
            unserialize(data[key])

def load_cache(mode):
    cache_loaded = get_state(f'{mode}_cache_loaded')
    if not (get_state(mode) or get_state(f'{mode}_load_started')):
        future = submit(load_json(get_path(mode), get_state('cache_queue'), mode))
        future.add_done_callback(lambda f: cache_loaded.set())
        set_state(f'{mode}_load_started', value=True)

async def cache_json(path: Path, data):
    data_copy = copy.deepcopy(data)
    serialize(data_copy)
    await asyncio.to_thread(_cache_json, path, data_copy)

def serialize(data: dict):
    for key in data:
        if isinstance(data[key], set):
            data[key] = list(data[key])
        elif isinstance(data[key], dict):
            serialize(data[key])

def _cache_json(path: Path, data):
    with open(path, 'w', encoding='utf-8') as f:
        json.dump(data, f, indent=4)

async def cache_txt(path: Path, data, queue: queue.Queue(), *attr):
    await asyncio.to_thread(_cache_txt, path, data)
    queue.put((*attr, True))

def _cache_txt(path: Path, data):
    with open(path, 'w', encoding='utf-8') as f:
        f.write(data)

def get_path(input_mode: str):
    return get_state('project_dir') / f'{input_mode}.json'

async def parse_uploaded_files(queue, uploaded, mode):
    for upload in uploaded:
        file_name = upload.name
        queue.put((file_name, mode, await asyncio.to_thread(parse_uploaded_file, upload)))

@st.dialog('Document Upload', width=1200)
def document_upload_dialog():

    shared_methods = ["JSON Upload", "CSV Upload", "SQL"]

    mode = "docs" if st.radio(
        "Input mode",
        ["Raw Documents", "Preprocessed Chunks"],
        horizontal=True,
        index=0
    ) == "Raw Documents" else "chunks"

    method = st.radio(
        "Input method",
        ["Manual File Upload"] + shared_methods if mode == "docs" else shared_methods,
        horizontal=True,
        index=0
    )

    q = get_state('upload_queue')

    if method == "Manual File Upload":
        uploaded = st.file_uploader(
            "Upload documents",
            accept_multiple_files=True,
            type=SUPPORTED_DOC_TYPES,
            key=f'{mode}_manual_upload'
        )

        if uploaded:
            names = [f.name for f in uploaded]
            name_counts = Counter(names)
            duplicates = [name for name, count in name_counts.items() if count > 1]

            if duplicates:
                st.toast(
                    f"Duplicate filenames detected in your upload: {', '.join(duplicates)}. "
                    "Please ensure each file has a unique name before uploading."
                )
            else:
                for f in uploaded:
                    if f.name in get_state('sources'):
                        st.markdown("A file of the same name already exists.")
                        col1, col2 = st.columns(2)
                        col1.button("Replace", type="secondary", key=f'{f.name}_duplicate_replace')
                        keep = col2.button("Keep Both", type="primary", key=f'{f.name}_duplicate_keep')

                        if keep:
                            num_similar_files = max({0} | {
                                int(m.group(1))
                                for s in get_state('sources')
                                if (m := re.search(rf'\(([0-9]+)\) {f.name}', s))
                            }) + 1
                            f.name = f"({num_similar_files}) {f.name}"
                        
                        set_state('sources', f.name, value=dict(), aggressive=False)
                        break
                    
                    set_state('sources', f.name, value=dict(), aggressive=False)

                get_state('coro_queue').put(('uploading', parse_uploaded_files(q, uploaded, mode)))
                set_state('active_dialog', value=None)
                st.rerun()

    elif method == "JSON Upload":
        f = st.file_uploader("Upload JSON file", type=["json"], key=f'{mode}_json_file')
        content_label = st.text_input("Content Label", key=f'{mode}_json_content_label')
        if st.button("Upload", key=f'{mode}_json_upload'):
            get_state('coro_queue').put(('uploading', _records_from_json(f, q, mode)))
            set_state('active_dialog', value=None)
            st.rerun()

    elif method == "CSV Upload":
        f = st.file_uploader("Upload CSV file", type=["csv"], key=f'{mode}_csv_upload')
        if f:
            get_state('coro_queue').put(('uploading', _records_from_csv(f, q, mode)))
            set_state('active_dialog', value=None)
            st.rerun()

    elif method == "SQL":
        _sql_form("raw_sql", q, mode)

def save_docs_or_chunks_to_cache(mode, key, content, chunking=False):
    parent_type = 'docs' if chunking else 'sources'
    existing = get_state(parent_type, key, mode, default=set()).copy()

    text_dir = get_state('project_dir') / mode
    text_dir.mkdir(exist_ok=True)

    text_ids = set()
    new_text_ids = set()
    
    for text in content:
        text_id = hash(text)
        text_ids.add(text_id)
        existing.discard(text_id)

        if 'upload_finished' not in get_state(mode, text_id, default=dict()):
            new_text_ids.add(text_id)

            get_state('coro_queue').put((
                'caching',
                cache_txt(
                    text_dir / f'{text_id}.txt',
                    text,
                    get_state('cache_queue'),
                    mode,
                    text_id,
                    'upload_finished'
                )
            ))
            
            set_state(mode, text_id, 'length', value=len(text))
        
    for id in existing:
        text_path = text_dir / f'{id}.txt'
        os.remove(text_path)
        delete_value(mode, id)
    
    if new_text_ids or existing:
        set_state(parent_type, key, mode, value=text_ids)
    
    return new_text_ids, existing

def manage_docs_and_chunks_in_memory():
    submit_coroutines()

    content_types = ['sources', 'docs', 'chunks']
    to_cache = set()

    all_set = True
    for content_type in content_types:
        all_set = all_set and get_state(f'{content_type}_cache_loaded').is_set()

    if not all_set:
        return
    
    cache_queue = get_state('cache_queue')

    while not cache_queue.empty():
        *attr, data = cache_queue.get_nowait()
        set_state(*attr, value=data)

        for content_type in content_types:
            if content_type in attr:
                to_cache.add(content_type)

    upload_queue = get_state('upload_queue')
    ignored_uploads = 0
    partial_uploads = 0
    
    while not upload_queue.empty():
        filename, mode, data = upload_queue.get_nowait()
        if not data:
            ignored_uploads += 1
        else:
            new_text_ids, deleted_ids = save_docs_or_chunks_to_cache(mode, filename, data['content'])
            
            if not new_text_ids:
                ignored_uploads += 1
            else:
                if len(new_text_ids) < len(data['content']):
                    partial_uploads += 1
                
                to_cache.add(mode)
                to_cache.add('sources')
            
            if deleted_ids:
                to_cache.add('sources')
    
    warnings = []
    upload_warnings = ""
    
    if ignored_uploads:
        upload_warnings += f"{ignored_uploads} upload(s) were ignored to avoid duplicate files. "
    if partial_uploads:
        upload_warnings += f"{partial_uploads} upload(s) partially contain preexisting documents. "
    
    if upload_warnings:
        upload_warnings += "Only unique documents were added."
        warnings.append(upload_warnings)
    
    chunking_errors = 0
    
    chunk_queue = get_state('chunk_queue')
    while not chunk_queue.empty():
        key, chunker_hash, chunk_texts = chunk_queue.get_nowait()
        if chunk_texts:
            new_chunk_ids, deleted_ids = save_docs_or_chunks_to_cache('chunks', key, chunk_texts, chunking=True)
            
            set_state('docs', 'key', 'chunking_strategy', value=chunker_hash)

            if new_chunk_ids:
                to_cache.add('chunks')
                to_cache.add('sources')
            if deleted_ids:
                to_cache.add('sources')

        else:
            chunking_errors += 1
    
    if chunking_errors:
        warnings.append(f"Unable to chunk {chunking_errors} documents. Only successful chunks were added.")
    
    delete_queue = get_state('delete_queue')
    while not delete_queue.empty():
        key = delete_queue.get_nowait()
        for doc_key in get_state('source', key, 'docs'):
            
            for chunk_key in get_state('docs', doc_key, 'chunks'):
                os.remove(get_state('project_dir') / 'chunks' / f'{chunk_key}.txt')
            
            delete_value('docs', doc_key)
            os.remove(get_state('project_dir') / 'docs' / f'{doc_key}.txt')
        
        for chunk_key in get_state('source', key, 'chunks'):
            os.remove(get_state('project_dir') / 'chunks' / f'{chunk_key}.txt')
        
        delete_value('sources', key)
    
    if warnings:
        st.toast(*warnings)

    for content_type in to_cache:
        get_state('coro_queue').put((
            'caching',
            cache_json(get_state('project_dir') / f'{content_type}.json', get_state(content_type))
        ))

@st.dialog("Delete Source", width=1200)
def delete_source(key):
    st.markdown(f"Are you sure you want to delete {key}?")
    col1, col2 = st.columns(2)

    if col1.button("Yes", type="secondary", key=f'{key}_delete_yes'):
        get_state('delete_queue').put(key)

    if col2.button("No", type="primary", key=f'{key}_delete_no'):
        set_state('active_dialog', value=None)

async def queue_load_txt(dir, display_queue):
    txt = await asyncio.to_thread(_load_txt, dir)
    display_queue.put(txt)

def content_display():
    display_queue = get_state('display_queue')
    while not display_queue.empty():
        set_state('display_text', value=display_queue.get_nowait())
    with st.container(height=400, border=True):
        st.markdown(get_state('display_text'))

@st.dialog("Source Inspector", width=1200, dismissible=False)
def source_inspector(key):
    docs = list(get_state('sources', key, 'docs'))
    doc_idx = st.number_input(
        "Document Index",
        value=0,
        step=1,
        min_value=0,
        max_value=len(docs) - 1
    )

    # Create a new queue object every time, so previous coroutines can't now modify it.
    if get_state('display_key') != f'{key}_{doc_idx}':
        set_state('display_key', value=f'{key}_{doc_idx}')
        set_state('display_queue', value=queue.Queue())
        display_queue = get_state('display_queue')

        get_state('coro_queue').put((
            'displaying',
            submit(queue_load_txt(
                get_state('project_dir') / 'docs' / f'{docs[doc_idx]}.txt',
                display_queue
            ))
        ))

    content_display()

    if st.button('Close', type='secondary', key=f'{key}_close'):
        set_state('active_source', value=None)
        set_state('display_key', value=None)
        set_state('display_text', value="Loading content...")
        set_state('display_queue', value=queue.Queue())
        close_dialog()

def sources_list():

    for key in get_state('sources').keys():
        with st.container(border=True):
            st.markdown(f"**{key}**")
            col1, col2 = st.columns(2)
            if col1.button("Open", key=f"open_{key}"):
                set_state('active_source', value=key)
                set_state('active_dialog', value='source_inspector')
            if col2.button("Delete", key=f"delete_{key}"):
                set_state('active_source', value=key)
                set_state('active_dialog', value='delete_source')

def document_metrics():
    raw_docs = get_state('docs').values()
    col1, col2 = st.columns(2)
    with col1:
        st.markdown(metric_card("DOCUMENTS", len(raw_docs)), unsafe_allow_html=True)
    with col2:
        if len(raw_docs) == 0:
            st.markdown(metric_card("AVG LENGTH (CHARS)", 0), unsafe_allow_html=True)
        else:
            st.markdown(metric_card(
                "AVG LENGTH (CHARS)",
                round(sum([doc['length'] for doc in raw_docs]) / len(raw_docs))
            ), unsafe_allow_html=True)

def chunk_metrics():
    chunks = get_state('chunks').values()
    col1, col2 = st.columns(2)
    with col1:
        st.markdown(metric_card("CHUNKS", len(chunks)), unsafe_allow_html=True)
    with col2:
        if len(chunks) == 0:
            st.markdown(metric_card("AVG LENGTH (CHARS)", 0), unsafe_allow_html=True)
        else:
            st.markdown(metric_card("AVG LENGTH (CHARS)", sum([chunk['length'] for chunk in chunks]) / len(chunks)), unsafe_allow_html=True)
    _chunk_inspector(get_state('chunks'))


@fragment(run_every=1)
def render_documents_tab():

    st.header("Documents")

    load_cache('sources')
    load_cache('docs')
    load_cache('chunks')
    
    manage_docs_and_chunks_in_memory()

    col1, _ = st.columns(2)
    if col1.button('Upload / Load Documents'):
        set_state('active_dialog', value='document_upload_dialog')
    
    col3, col4 = st.columns([3, 1])
    with col4:
        with st.container(height=400, border=False):
            sources_list()
    
    with col3:
        document_metrics()

        if st.button("Chunk Documents"):
            set_state('active_dialog', value='chunking_dialog')
        
        chunk_metrics()
    
    render_nlp_analytics_tab()
    
    if get_state('active_dialog') == 'document_upload_dialog':
        document_upload_dialog()
    elif get_state('active_dialog') == 'chunking_dialog':
        chunking_dialog(get_state('docs'))
    elif get_state('active_dialog') == 'source_inspector':
        source_inspector(get_state('active_source'))
    elif get_state('active_dialog') == 'delete_source':
        delete_source(get_state('active_source'))


# ---------------------------------------------------------------------------
# Dev entrypoint
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    st.set_page_config(page_title="Documents", layout="wide")
    render_documents_tab()